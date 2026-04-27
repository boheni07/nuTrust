from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import copy

# ── 색상 정의 ──────────────────────────────────────
BLUE       = RGBColor(0x25, 0x63, 0xEB)   # Trust Blue
DARK_BLUE  = RGBColor(0x1E, 0x40, 0xAF)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x0F, 0x17, 0x2A)
GRAY       = RGBColor(0x64, 0x74, 0x8B)
BG         = RGBColor(0xF8, 0xFA, 0xFC)
GREEN      = RGBColor(0x05, 0x96, 0x69)
RED        = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ── 헬퍼 함수 ──────────────────────────────────────
def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background() if line is None else None
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=DARK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_multiline(slide, lines, left, top, width, height,
                  size=15, color=DARK, bold_first=False, line_spacing=None):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if line_spacing:
            p.space_after = Pt(line_spacing)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = (bold_first and i == 0)
    return txb

def page_bg(slide, color=BG):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=color)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill=BLUE)
    add_text(slide, title, 0.5, 0.15, 12, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.5, 0.8, 12, 0.4, size=14, color=LIGHT_BLUE)

def slide_number(slide, num, total=16):
    add_text(slide, f"{num} / {total}", 12.3, 7.1, 1, 0.3, size=11, color=GRAY, align=PP_ALIGN.RIGHT)

def draw_table(slide, headers, rows, left, top, width, height,
               col_widths=None, header_size=13, row_size=12):
    n_cols = len(headers)
    if col_widths is None:
        col_widths = [width / n_cols] * n_cols

    row_h = height / (len(rows) + 1)
    x = left

    # 헤더
    for i, (h, cw) in enumerate(zip(headers, col_widths)):
        add_rect(slide, x, top, cw - 0.02, row_h - 0.02, fill=BLUE)
        add_text(slide, h, x + 0.05, top + 0.05, cw - 0.1, row_h - 0.1,
                 size=header_size, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += cw

    # 데이터 행
    for r_idx, row in enumerate(rows):
        x = left
        bg = LIGHT_BLUE if r_idx % 2 == 0 else WHITE
        for i, (cell, cw) in enumerate(zip(row, col_widths)):
            add_rect(slide, x, top + row_h * (r_idx + 1), cw - 0.02, row_h - 0.02, fill=bg)
            txt_color = GREEN if str(cell).startswith('✅') else (RED if str(cell).startswith('❌') else DARK)
            add_text(slide, str(cell), x + 0.05,
                     top + row_h * (r_idx + 1) + 0.05,
                     cw - 0.1, row_h - 0.1,
                     size=row_size, color=txt_color, align=PP_ALIGN.CENTER, wrap=True)
            x += cw


# ══════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=BLUE)
add_rect(s, 0, 2.8, 13.33, 0.04, fill=WHITE)

add_text(s, "nuTrust ServiceDesk", 1, 1.2, 11.33, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "티켓 기반 요청 관리를 통한 고객 신뢰 구축 플랫폼",
         1, 2.4, 11.33, 0.7, size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

info = [
    ("기간",        "2026-04-07 ~ 2026-04-08"),
    ("방법론",       "PDCA  (PM → Plan → Design → Do → Check → Act)"),
    ("최종 Match Rate", "~95%  (서비스)  /  100%  (도메인)"),
    ("성공 기준",    "6 / 6  (100%)  달성"),
]
for i, (k, v) in enumerate(info):
    y = 3.3 + i * 0.65
    add_text(s, k, 2.5, y, 2.5, 0.5, size=14, bold=True, color=LIGHT_BLUE)
    add_text(s, v, 5.0, y, 6.5, 0.5, size=14, color=WHITE)


# ══════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "목차")
slide_number(s, 2)

items = [
    ("01", "문제 정의",    "왜 만들었나 — 고객 페인포인트 & 시장 공백"),
    ("02", "솔루션",       "무엇을 만들었나 — 3역할 플랫폼 & 핵심 가치"),
    ("03", "PDCA 여정",    "어떻게 만들었나 — PM → Plan → Design → Do → Check → Act"),
    ("04", "구현 결과",    "얼마나 만들었나 — 95파일 / 46 API / 16 UI"),
    ("05", "품질 검증",    "얼마나 잘 만들었나 — Match Rate 86% → 95% → 100%"),
    ("06", "교훈 & 스택",  "무엇을 배웠나 — 8가지 교훈 + 기술 스택"),
]
for i, (num, title, desc) in enumerate(items):
    y = 1.5 + i * 0.9
    add_rect(s, 0.6, y, 0.6, 0.6, fill=BLUE)
    add_text(s, num, 0.6, y + 0.08, 0.6, 0.45, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.4, y + 0.02, 2.5, 0.35, size=17, bold=True, color=DARK_BLUE)
    add_text(s, desc,  1.4, y + 0.35, 10, 0.35, size=13, color=GRAY)


# ══════════════════════════════════════════════════
# 슬라이드 3 — 문제 정의
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "01. 문제 정의", "고객 페인포인트 & 시장 공백")
slide_number(s, 3)

add_text(s, "고객 페인포인트 3가지", 0.5, 1.4, 8, 0.4, size=16, bold=True, color=DARK_BLUE)
draw_table(s,
    ["#", "문제", "영향"],
    [
        ["P1", "요청 추적 불가 — 이메일/전화/메신저 분산", "60~70% 재문의 발생"],
        ["P2", "처리 과정 불투명 — '블랙박스' 경험",     "NPS 20~30점 하락"],
        ["P3", "완료 기준 불일치 — 일방적 완료 처리",    "재문의율 25~35%"],
    ],
    left=0.5, top=1.85, width=12.33, height=1.7,
    col_widths=[0.6, 7.0, 4.73]
)

add_text(s, "시장 공백 분석", 0.5, 3.8, 8, 0.4, size=16, bold=True, color=DARK_BLUE)
draw_table(s,
    ["경쟁사", "문제점", "nuTrust 차별점"],
    [
        ["ServiceNow",       "도입 3~6개월, 과도하게 비쌈",          ""],
        ["Zendesk/Freshdesk","고객 승인 워크플로우 약함, 한국어 부족", ""],
        ["✅ nuTrust",       "적정 가격 + 고객 승인 기반 차별화",     "50~500명 한국 중견기업"],
    ],
    left=0.5, top=4.25, width=12.33, height=2.0,
    col_widths=[2.5, 5.5, 4.33]
)


# ══════════════════════════════════════════════════
# 슬라이드 4 — 솔루션 개요
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "02. 솔루션 개요", "고객 승인 기반 서비스 완결성")
slide_number(s, 4)

add_rect(s, 0.5, 1.4, 12.33, 0.8, fill=LIGHT_BLUE)
add_text(s, "핵심 가치: 고객이 직접 결과를 확인·승인하는 구조로 신뢰를 체계적으로 구축",
         0.6, 1.5, 12, 0.6, size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

roles = [
    ("고객 포털",      ["티켓 등록", "실시간 상태 추적", "결과 승인/반려", "CSAT 평가"]),
    ("담당자 대시보드", ["배정 티켓 처리", "처리계획 등록", "내부메모 소통", "상태 전환"]),
    ("관리자 콘솔",    ["SLA 모니터링", "팀 성과 분석", "프로젝트 관리", "보고서"]),
]
colors = [BLUE, DARK_BLUE, RGBColor(0x03, 0x69, 0x9E)]
for i, (role, items) in enumerate(roles):
    x = 0.5 + i * 4.2
    add_rect(s, x, 2.4, 3.9, 0.55, fill=colors[i])
    add_text(s, role, x, 2.45, 3.9, 0.5, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 2.95, 3.9, 2.3, fill=WHITE)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", x + 0.15, 3.0 + j * 0.52, 3.7, 0.5, size=13, color=DARK)

add_text(s, "성공 기준 (SC)", 0.5, 5.5, 3, 0.4, size=15, bold=True, color=DARK_BLUE)
scs = ["CSAT ≥ 4.0/5.0", "SLA ≥ 95%", "승인율 ≥ 85%", "재문의율 ≤ 8%"]
for i, sc in enumerate(scs):
    x = 0.5 + i * 3.1
    add_rect(s, x, 5.95, 2.9, 0.55, fill=BLUE)
    add_text(s, sc, x, 5.98, 2.9, 0.5, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# 슬라이드 5 — PDCA 여정
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "03. PDCA 여정", "2일간의 전체 개발 플로우")
slide_number(s, 5)

phases = [
    ("PM",     "PRD 작성\n4개 에이전트\n병렬 분석"),
    ("PLAN",   "SC 6개 확정\nMVP 범위\n합의"),
    ("DESIGN", "아키텍처 3안\nOption C\n선택"),
    ("DO",     "8 모듈 구현\n95파일\n~3,000 LOC"),
    ("CHECK",  "Gap 분석\n86%→91%\n→95%"),
    ("ACT",    "수정 2회\nSC 6/6 ✅\n15/15 PASS"),
]
pw = 1.9
for i, (phase, desc) in enumerate(phases):
    x = 0.4 + i * 2.08
    add_rect(s, x, 1.5, pw, 0.65, fill=BLUE)
    add_text(s, phase, x, 1.55, pw, 0.55, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, 2.15, pw, 2.0, fill=WHITE)
    add_text(s, desc, x + 0.05, 2.2, pw - 0.1, 1.9, size=13, color=DARK, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, "→", x + pw - 0.05, 1.88, 0.4, 0.4, size=18, bold=True, color=BLUE)

add_text(s, "총 기간: 2일  (2026-04-07 ~ 2026-04-08)", 0.4, 4.4, 12.33, 0.5,
         size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

draw_table(s,
    ["날짜", "세션", "결과"],
    [
        ["04-07 AM", "PM + Plan", "PRD + Plan 문서 확정"],
        ["04-07 AM", "Design",   "Option C + Design Anchor 확정"],
        ["04-07 PM", "Do",       "95파일 / 46 API / 16 UI 완성"],
        ["04-07~08", "Check+Act","86% → 91.2% → 95%, SC 6/6"],
        ["04-08",    "ticket-machine 점검", "G-01 수정, 100% Match Rate"],
    ],
    left=0.5, top=5.0, width=12.33, height=2.2,
    col_widths=[2.0, 3.5, 6.83], row_size=11
)


# ══════════════════════════════════════════════════
# 슬라이드 6 — PM & Plan
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 1 — PM & Plan", "PRD 자동 생성 + SC 확정")
slide_number(s, 6)

add_text(s, "PM 에이전트 팀 병렬 분석", 0.5, 1.4, 6, 0.4, size=15, bold=True, color=DARK_BLUE)
agents = [
    ("pm-lead",      "오케스트레이터",  "4개 에이전트 조율"),
    ("pm-discovery", "기회 발견",      "OST (Opportunity Solution Tree)"),
    ("pm-strategy",  "전략 분석",      "JTBD 6-Part + Lean Canvas"),
    ("pm-research",  "시장 조사",      "3 페르소나 + 5 경쟁사 + TAM/SAM/SOM"),
    ("pm-prd",       "PRD 종합",       "8섹션 PRD 자동 생성"),
]
for i, (name, role, desc) in enumerate(agents):
    y = 1.9 + i * 0.62
    c = BLUE if i == 0 else DARK_BLUE
    add_rect(s, 0.5, y, 2.0, 0.5, fill=c)
    add_text(s, name, 0.5, y + 0.07, 2.0, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, f"{role} — {desc}", 2.65, y + 0.1, 3.8, 0.4, size=12, color=DARK)

add_text(s, "Success Criteria 6개 확정", 7.0, 1.4, 6, 0.4, size=15, bold=True, color=DARK_BLUE)
draw_table(s,
    ["#", "기준", "구현 방식"],
    [
        ["SC-1", "상태 머신 전체 흐름", "XState 5, 9 states"],
        ["SC-2", "SLA 4근무시간 자동 접수", "비즈니스 캘린더 + Cron"],
        ["SC-3", "연기 3중 가드", "횟수 / 상태 / 날짜"],
        ["SC-4", "파일 첨부", "S3 presigned URL"],
        ["SC-5", "RBAC 3역할", "미들웨어 기반"],
        ["SC-6", "이메일 알림", "7개 템플릿"],
    ],
    left=7.0, top=1.85, width=5.8, height=5.2,
    col_widths=[0.7, 2.4, 2.7], row_size=11
)


# ══════════════════════════════════════════════════
# 슬라이드 7 — Design
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 2 — Design", "아키텍처 3안 비교 → Option C 선택")
slide_number(s, 7)

draw_table(s,
    ["", "A: Monolithic", "B: Clean/DDD", "✅ C: Feature-based"],
    [
        ["파일 수",  "~40",    "~80+",  "~55"],
        ["복잡도",  "Low",    "High",  "Medium"],
        ["유지보수", "Medium", "High",  "High"],
        ["MVP 속도", "Fast",   "Slow",  "Balanced"],
    ],
    left=0.5, top=1.4, width=8.5, height=2.8,
    col_widths=[1.5, 2.0, 2.0, 3.0]
)
add_rect(s, 0.5, 4.35, 8.5, 0.55, fill=LIGHT_BLUE)
add_text(s, "선택 이유: MVP 속도와 장기 유지보수의 균형", 0.6, 4.4, 8.3, 0.45,
         size=13, bold=True, color=DARK_BLUE)

add_text(s, "Design Anchor 확정", 9.2, 1.4, 4, 0.4, size=15, bold=True, color=DARK_BLUE)
anchor_items = [
    ("Primary Color", "#2563EB  Trust Blue"),
    ("Typography", "Pretendard (KO) + Inter (EN)"),
    ("Layout", "Sidebar(256px) + Content"),
    ("Card Radius", "rounded-lg (8px)"),
    ("Tone", "전문적·신뢰감"),
]
for i, (k, v) in enumerate(anchor_items):
    y = 1.9 + i * 0.65
    add_rect(s, 9.2, y, 1.5, 0.52, fill=BLUE)
    add_text(s, k, 9.22, y + 0.08, 1.46, 0.38, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, v, 10.8, y + 0.1, 2.4, 0.38, size=11, color=DARK)

add_text(s, "핵심 설계 결정", 0.5, 5.1, 8, 0.4, size=15, bold=True, color=DARK_BLUE)
decisions = [
    "XState 5 + 서버 사이드 validateTransition() 이중 구조",
    "9-state 티켓 머신 (REGISTERED → CLOSED)",
    "연기 3중 가드 (횟수 / 상태 / 날짜)",
    "NextAuth.js v5 미들웨어 기반 RBAC",
]
for i, d in enumerate(decisions):
    add_text(s, f"• {d}", 0.5, 5.55 + i * 0.45, 12, 0.42, size=13, color=DARK)


# ══════════════════════════════════════════════════
# 슬라이드 8 — Do: 구현
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 3 — Do: 구현", "8 모듈 순차 구현")
slide_number(s, 8)

draw_table(s,
    ["모듈", "내용", "주요 파일"],
    [
        ["module-1", "인프라 + 도메인",       "schema.prisma, machine.ts, auth.ts"],
        ["module-2", "티켓 CRUD",             "tickets/route.ts, [id]/route.ts"],
        ["module-3", "상태 전환 API",          "accept, start, complete, assign"],
        ["module-4", "승인/연기 API",          "approve, reject, postpone"],
        ["module-5", "소통/파일/CSAT",         "comments, history, upload, csat"],
        ["module-6", "관리자 API",             "clients, projects, teams, sla"],
        ["module-7", "대시보드/크론",           "dashboard, reports, cron"],
        ["module-8", "UI 16페이지",            "portal / agent / admin 레이아웃"],
    ],
    left=0.5, top=1.4, width=8.0, height=5.8,
    col_widths=[1.5, 2.7, 3.8], row_size=11
)

add_text(s, "구현 중 해결한 이슈", 8.7, 1.4, 4.5, 0.4, size=15, bold=True, color=DARK_BLUE)
issues = [
    ("Prisma 7 Breaking Change", "v6 다운그레이드"),
    ("XState 5 타입 오류", "정적 context 객체 패턴"),
    ("Route Group 경로 충돌", "role prefix 추가"),
    ("접수 2단계 UX 문제", "ACCEPT+START 통합 API"),
]
for i, (prob, sol) in enumerate(issues):
    y = 1.9 + i * 1.15
    add_rect(s, 8.7, y, 4.4, 0.5, fill=RGBColor(0xFE, 0xF2, 0xF2))
    add_text(s, f"문제: {prob}", 8.8, y + 0.02, 4.2, 0.25, size=11, bold=True, color=RED)
    add_rect(s, 8.7, y + 0.5, 4.4, 0.45, fill=RGBColor(0xEC, 0xFD, 0xF5))
    add_text(s, f"해결: {sol}", 8.8, y + 0.52, 4.2, 0.25, size=11, bold=True, color=GREEN)


# ══════════════════════════════════════════════════
# 슬라이드 9 — 상태 머신
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "티켓 상태 머신", "XState 5 기반 9-State 설계")
slide_number(s, 9)

states = [
    ("REGISTERED",            1.5,  1.5),
    ("ACCEPTED",              4.5,  1.5),
    ("IN_PROGRESS",           7.5,  1.5),
    ("DELAYED",               7.5,  3.2),
    ("COMPLETION_REQUESTED",  7.5,  4.9),
    ("POSTPONEMENT_REQUESTED",10.5, 1.5),
    ("APPROVED",              10.5, 4.9),
    ("CLOSED",                10.5, 6.4),
]
sw, sh = 2.3, 0.6
for name, x, y in states:
    is_final = name in ("CLOSED",)
    fill = GREEN if is_final else BLUE
    add_rect(s, x, y, sw, sh, fill=fill)
    add_text(s, name, x, y + 0.08, sw, sh - 0.15, size=10, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)

transitions = [
    ("ACCEPT / AUTO_ACCEPT", 3.8, 1.75),
    ("START",                6.8, 1.75),
    ("DELAY",                7.5, 2.4),
    ("REQUEST_COMPLETION",   7.5, 4.3),
    ("REQUEST_POSTPONEMENT\n(3-guard)", 9.3, 1.75),
    ("APPROVE",              10.5, 4.15),
    ("REJECT → IN_PROGRESS", 9.2, 4.5),
    ("APPROVE_POSTPONEMENT\n→ IN_PROGRESS", 9.0, 2.8),
    ("SUBMIT_CSAT",          10.9, 5.8),
]
for label, x, y in transitions:
    add_text(s, label, x, y, 1.8, 0.6, size=8, color=GRAY, italic=True)

add_rect(s, 0.5, 6.5, 12.33, 0.65, fill=LIGHT_BLUE)
add_text(s, "3-guard 연기 조건:  ① postponementCount < 1  ②  status ≠ DELAYED  ③  now < plannedDueDate",
         0.6, 6.56, 12, 0.5, size=12, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# 슬라이드 10 — Check & Act
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 4 — Check & Act", "Gap 분석 → 2회 반복 개선")
slide_number(s, 10)

draw_table(s,
    ["라운드", "Overall", "Structural", "Functional", "Contract"],
    [
        ["Check-1",  "86%",   "88%", "82%", "89%"],
        ["Act-1",    "91.2%", "90%", "90%", "93%"],
        ["Act-2 ✅", "~95%",  "95%", "95%", "95%"],
    ],
    left=0.5, top=1.4, width=6.5, height=2.2,
    col_widths=[1.5, 1.3, 1.3, 1.3, 1.1]
)

add_text(s, "Act-1: 사용자 피드백 기반 수정 (5건)", 7.2, 1.4, 5.9, 0.4, size=14, bold=True, color=DARK_BLUE)
act1 = [
    ("I-1", '"댓글 입력은 필수예요"', "TicketDetail 댓글 UI 추가"),
    ("I-2", '"첨부파일 선택도요"',    "TicketForm 파일 업로드 추가"),
    ("I-3", '"이건 못 쓰겠는데요"',   "API 기반 드롭다운 교체"),
    ("I-4/5", "PUT API 누락",         "users/:id, sla-policies/:id route 추가"),
]
for i, (num, feedback, fix) in enumerate(act1):
    y = 1.9 + i * 0.65
    add_rect(s, 7.2, y, 0.55, 0.52, fill=RED)
    add_text(s, num, 7.2, y + 0.1, 0.55, 0.35, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, feedback, 7.85, y + 0.01, 2.8, 0.28, size=10, italic=True, color=GRAY)
    add_text(s, f"→ {fix}", 7.85, y + 0.28, 5.2, 0.28, size=10, color=DARK)

add_text(s, "Act-2: 추가 요청 반영 (8건)", 0.5, 3.85, 6.5, 0.4, size=14, bold=True, color=DARK_BLUE)
act2 = [
    "비인증 시 빈 화면 → middleware.ts 리다이렉트",
    "SLA 카운트다운 배지 (초과/임박/정상)",
    "프로젝트/기간 필터",
    "관리자/담당자 설정 페이지",
    "approval / csat / clients / projects 모듈 분리",
]
for i, item in enumerate(act2):
    add_text(s, f"• {item}", 0.5, 4.35 + i * 0.48, 12, 0.45, size=13, color=DARK)

add_rect(s, 0.5, 6.85, 12.33, 0.5, fill=BLUE)
add_text(s, "최종: SC 6/6  ✅  |  API 테스트 15/15 PASSED  |  Match Rate ~95%",
         0.6, 6.88, 12, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════
# 슬라이드 11 — ticket-machine 점검
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 5 — ticket-machine 점검", "도메인 레이어 독립 PDCA")
slide_number(s, 11)

draw_table(s,
    ["Axis", "Rate", "Notes"],
    [
        ["Structural",   "100%",          "4/4 파일 (machine, guards, types, index)"],
        ["Functional",   "97% → 100%",    "G-01 수정 후"],
        ["Contract",     "100%",          "오류 코드 4개 완전 일치"],
        ["Overall ✅",   "100%",          "G-01 수정 완료"],
    ],
    left=0.5, top=1.4, width=7.5, height=2.3,
    col_widths=[2.0, 1.8, 3.7]
)

add_text(s, "발견 이슈 G-01 (Minor)", 0.5, 3.9, 7.5, 0.4, size=15, bold=True, color=DARK_BLUE)

add_rect(s, 0.5, 4.35, 5.8, 0.5, fill=RGBColor(0xFE, 0xF2, 0xF2))
add_text(s, "Before — 도달 불가 dead state", 0.6, 4.37, 5.6, 0.25, size=11, bold=True, color=RED)
before_code = "REJECTED: { type: 'final' }   ← 어떤 전환도 이 state로 오지 않음"
add_text(s, before_code, 0.6, 4.62, 5.6, 0.3, size=10, color=DARK)

add_rect(s, 0.5, 5.0, 5.8, 0.5, fill=RGBColor(0xEC, 0xFD, 0xF5))
add_text(s, "After — 제거 (types.ts / schema.prisma는 유지)", 0.6, 5.02, 5.6, 0.25, size=11, bold=True, color=GREEN)
after_code = "REJECTED state 삭제 → XState 8개 도달 가능 상태만 유지"
add_text(s, after_code, 0.6, 5.27, 5.6, 0.3, size=10, color=DARK)

add_text(s, "guards.ts 구현 목록", 6.5, 1.4, 6.3, 0.4, size=15, bold=True, color=DARK_BLUE)
guards = [
    "canRequestPostponement()  — 3중 조건 검증",
    "canRequestCompletion()    — IN_PROGRESS / DELAYED",
    "canAccept()               — REGISTERED 상태",
    "canApproveOrReject()      — COMPLETION_REQUESTED",
    "canRespondToPostponement()— POSTPONEMENT_REQUESTED",
    "getPostponementRejectionReason()  — 상세 오류 메시지",
]
for i, g in enumerate(guards):
    add_text(s, f"• {g}", 6.5, 1.9 + i * 0.55, 6.3, 0.52, size=12, color=DARK)


# ══════════════════════════════════════════════════
# 슬라이드 12 — 환경설정 & 실행
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "Session 6 — 환경설정 & 서비스 실행", "Docker PostgreSQL + 시드 → 정상 실행")
slide_number(s, 12)

add_text(s, "문제 → 해결 흐름", 0.5, 1.4, 6, 0.4, size=15, bold=True, color=DARK_BLUE)
steps = [
    (RED,   "로그인 실패"),
    (GRAY,  "DB 연결 오류 확인 (P1000)"),
    (GRAY,  "PostgreSQL 미실행 확인 (port 5432 없음)"),
    (BLUE,  "Docker PostgreSQL 컨테이너 실행"),
    (BLUE,  "prisma db push (스키마 적용)"),
    (BLUE,  "npx tsx prisma/seed.ts (시드 데이터)"),
    (GREEN, "npm run dev → http://localhost:3000 정상 ✅"),
]
for i, (c, text) in enumerate(steps):
    y = 1.9 + i * 0.63
    add_rect(s, 0.5, y, 0.08, 0.45, fill=c)
    add_text(s, text, 0.75, y + 0.05, 5.5, 0.4, size=13, color=DARK)
    if i < 6:
        add_text(s, "↓", 0.5, y + 0.45, 0.4, 0.3, size=12, color=GRAY)

add_text(s, "시드 데이터 & 로그인 계정", 6.8, 1.4, 6.2, 0.4, size=15, bold=True, color=DARK_BLUE)
draw_table(s,
    ["역할", "이메일", "비밀번호"],
    [
        ["관리자", "admin@servicedesk.com",   "password123"],
        ["매니저", "manager@servicedesk.com", "password123"],
        ["담당자", "agent1@servicedesk.com",  "password123"],
        ["고객",   "customer1@abc.com",        "password123"],
    ],
    left=6.8, top=1.9, width=6.0, height=2.4,
    col_widths=[1.4, 3.0, 1.6]
)

add_text(s, "시드 결과", 6.8, 4.5, 6, 0.4, size=14, bold=True, color=DARK_BLUE)
seed_results = [
    "2개 고객사 (ABC제조, XYZ금융)",
    "5개 부서, 3개 담당 연락처",
    "4개 프로젝트, 5개 배정",
    "8명 사용자, 3개 SLA 정책",
    "3개 샘플 티켓 (IN_PROGRESS / REGISTERED / COMPLETION_REQUESTED)",
]
for i, r in enumerate(seed_results):
    add_text(s, f"• {r}", 6.8, 4.95 + i * 0.48, 6.0, 0.45, size=12, color=DARK)


# ══════════════════════════════════════════════════
# 슬라이드 13 — 최종 산출물
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "최종 산출물", "구현 규모 & SC 달성 현황")
slide_number(s, 13)

metrics = [
    ("95",    "Source Files"),
    ("~3,000","Lines of Code"),
    ("46",    "API Endpoints"),
    ("16",    "UI Pages"),
    ("14",    "DB Models"),
    ("6",     "Feature Modules"),
]
mw = 2.0
for i, (val, label) in enumerate(metrics):
    x = 0.5 + i * 2.08
    add_rect(s, x, 1.4, mw, 1.1, fill=BLUE)
    add_text(s, val, x, 1.45, mw, 0.65, size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, label, x, 2.1, mw, 0.35, size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(s, "Success Criteria 최종 달성", 0.5, 2.75, 12, 0.4, size=15, bold=True, color=DARK_BLUE)
draw_table(s,
    ["#", "Success Criteria", "상태", "증거"],
    [
        ["SC-1", "상태 머신 전체 흐름",    "✅ Met", "XState 9 states + 15/15 API PASS"],
        ["SC-2", "SLA 4근무시간 자동 접수", "✅ Met", "businessCalendar.ts + cron/route.ts"],
        ["SC-3", "연기 3중 가드",           "✅ Met", "guards.ts + 3개 상세 오류 코드"],
        ["SC-4", "파일 첨부",               "✅ Met", "S3 presign + 댓글 첨부 + 다운로드"],
        ["SC-5", "RBAC 3역할",              "✅ Met", "middleware.ts + requireAuth 전 API"],
        ["SC-6", "이메일 알림",             "✅ Met", "7개 템플릿 (dev console fallback)"],
    ],
    left=0.5, top=3.2, width=12.33, height=3.8,
    col_widths=[0.7, 2.8, 1.0, 7.83], row_size=12
)


# ══════════════════════════════════════════════════
# 슬라이드 14 — 교훈 & 기술 스택
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
page_bg(s)
header_bar(s, "주요 교훈 & 기술 스택", "8가지 교훈 + 기술 스택")
slide_number(s, 14)

add_text(s, "주요 교훈", 0.5, 1.4, 6.5, 0.4, size=15, bold=True, color=DARK_BLUE)
lessons = [
    ("1", "Next.js Route Group 충돌",  "role prefix 필수 (/agent-dashboard)"),
    ("2", "Prisma 메이저 버전 주의",   "v6 고정 사용"),
    ("3", "XState 5 Context 타입",     "정적 context 객체 패턴"),
    ("4", "미들웨어는 필수",           "페이지 레벨 인증 필요"),
    ("5", "ACCEPT+START 통합",         "2단계 → 1 API로 UX 개선"),
    ("6", "Feature 모듈 분리 타이밍",  "MVP 후 안정화 시점에 분리"),
    ("7", "Dead State 관리",           "Check 단계에서 발견·제거"),
    ("8", "DB 환경 코드화",            "Docker Compose 권장"),
]
for i, (num, prob, sol) in enumerate(lessons):
    y = 1.9 + i * 0.63
    add_rect(s, 0.5, y, 0.45, 0.5, fill=BLUE)
    add_text(s, num, 0.5, y + 0.1, 0.45, 0.35, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, prob, 1.05, y + 0.01, 2.6, 0.27, size=11, bold=True, color=DARK)
    add_text(s, f"→ {sol}", 1.05, y + 0.27, 5.8, 0.27, size=11, color=GRAY)

add_text(s, "기술 스택", 7.3, 1.4, 5.8, 0.4, size=15, bold=True, color=DARK_BLUE)
draw_table(s,
    ["Layer", "Technology", "Version"],
    [
        ["Framework",   "Next.js App Router",    "16.2.2"],
        ["Database",    "PostgreSQL + Prisma",   "16 + 6.19"],
        ["State Machine","XState",               "5.30"],
        ["UI",          "Tailwind + shadcn/ui",  "4.x"],
        ["Auth",        "NextAuth.js",           "v5 beta"],
        ["Client State","Zustand + TanStack Q",  "5.x"],
        ["Forms",       "react-hook-form + Zod", "7.x + 4.x"],
        ["Container",   "Docker",                "29.2.1"],
    ],
    left=7.3, top=1.9, width=5.7, height=5.3,
    col_widths=[2.0, 2.3, 1.4], row_size=11
)


# ══════════════════════════════════════════════════
# 슬라이드 15 — 마무리
# ══════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=BLUE)
add_rect(s, 0, 3.1, 13.33, 0.04, fill=WHITE)

add_text(s, "프로젝트 완료", 1, 0.8, 11.33, 0.8,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "nuTrust ServiceDesk", 1, 1.55, 11.33, 0.65,
         size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

summary = [
    ("Match Rate",    "~95% (서비스)  /  100% (도메인)"),
    ("Success Criteria", "6 / 6  (100%)"),
    ("기간",          "2026-04-07 ~ 2026-04-08  (2일)"),
    ("산출물",        "95파일  /  ~3,000 LOC  /  46 API  /  16 UI"),
    ("핵심 가치",     "고객 승인 기반 서비스 완결성"),
]
for i, (k, v) in enumerate(summary):
    y = 3.4 + i * 0.68
    add_text(s, k, 2.0, y, 2.8, 0.55, size=14, bold=True, color=LIGHT_BLUE)
    add_text(s, v, 4.8, y, 7.5, 0.55, size=14, color=WHITE)

add_rect(s, 2.0, 6.8, 9.33, 0.5, fill=RGBColor(0x1E, 0x40, 0xAF))
add_text(s, "고객이 승인하는 서비스가 신뢰를 만든다",
         2.0, 6.82, 9.33, 0.42, size=15, bold=True, italic=True,
         color=WHITE, align=PP_ALIGN.CENTER)


# ── 저장 ───────────────────────────────────────────
output_path = r"E:\Google Drive\Dev\nuTrust\docs\04-report\nutrust-servicedesk.pptx"
prs.save(output_path)
print(f"DONE: {output_path}")
print(f"Slides: {len(prs.slides)}")
